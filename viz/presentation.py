"""
Presentation generator — turns analysis results into an editable PowerPoint deck.

Design goals (what makes this sellable, not a gimmick):
  - Real, editable .pptx (python-pptx) — the user rewrites wording if they want.
  - Charts are deterministic raster of real data (vl-convert) — never image-gen.
  - Template-loadable: pass a customer .pptx/.potx and the deck inherits their brand
    (theme colors, fonts, master). "Upload your template once, every deck on-brand."
  - Audit-trail appendix: the SQL behind every number — a board deck you can DEFEND.
    No competitor's auto-deck carries its own receipts.

A "deck" is a list of slide dicts the orchestrator/brain produces:
  {"type": "title", "title": ..., "subtitle": ...}
  {"type": "summary", "title": ..., "bullets": [...]}
  {"type": "chart", "title": ..., "takeaway": ..., "chart": <vega spec>}
  {"type": "appendix", "items": [{"title":..., "sql":...}, ...]}
"""
from __future__ import annotations

import io
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from viz.raster import vegalite_to_png

# Default theme (Teal Trust — matches the product brand). Overridden by a template.
@dataclass
class Theme:
    bg_dark: str = "0B1620"
    bg_light: str = "FFFFFF"
    ink: str = "1A2A35"
    ink_dim: str = "5A6B78"
    signal: str = "0F9B8E"
    accent: str = "34D6C4"
    title_font: str = "Cambria"     # safe-list serif w/ personality
    body_font: str = "Calibri"      # safe-list sans


def _rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h)


class DeckBuilder:
    def __init__(self, template_path: str | None = None, theme: Theme | None = None) -> None:
        # If a customer template is given, inherit its masters/branding.
        self.prs = Presentation(template_path) if template_path else Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = theme or Theme()
        self._blank = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[-1]

    # -- helpers -------------------------------------------------------------
    def _add(self):
        return self.prs.slides.add_slide(self._blank)

    def _box(self, slide, l, t, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        return tf

    def _bg(self, slide, hexc):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(hexc)

    def _para(self, tf, text, size, color, bold=False, font=None, align=PP_ALIGN.LEFT, first=False):
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font or self.theme.body_font
        r.font.color.rgb = _rgb(color)
        return p

    # -- slide types ---------------------------------------------------------
    def title_slide(self, title: str, subtitle: str = "") -> None:
        s = self._add()
        self._bg(s, self.theme.bg_dark)
        tf = self._box(s, 0.9, 2.5, 11.5, 2.4)
        self._para(tf, title, 44, "FFFFFF", bold=True, font=self.theme.title_font, first=True)
        if subtitle:
            self._para(tf, subtitle, 20, self.theme.accent, font=self.theme.body_font)

    def summary_slide(self, title: str, bullets: list[str]) -> None:
        s = self._add()
        self._bg(s, self.theme.bg_light)
        self._para(self._box(s, 0.8, 0.6, 11.7, 1.0), title, 32, self.theme.ink,
                   bold=True, font=self.theme.title_font, first=True)
        tf = self._box(s, 0.9, 1.9, 11.5, 5.0)
        for i, b in enumerate(bullets):
            p = self._para(tf, "•  " + b, 17, self.theme.ink, font=self.theme.body_font,
                           first=(i == 0))
            p.space_after = Pt(12)

    def chart_slide(self, title: str, takeaway: str, chart_spec: dict) -> None:
        s = self._add()
        self._bg(s, self.theme.bg_light)
        self._para(self._box(s, 0.8, 0.55, 11.7, 0.9), title, 28, self.theme.ink,
                   bold=True, font=self.theme.title_font, first=True)
        # takeaway (left), chart (right) — two-column, no accent bars
        if takeaway:
            tf = self._box(s, 0.8, 1.8, 4.3, 4.8)
            self._para(tf, takeaway, 16, self.theme.ink_dim, font=self.theme.body_font, first=True)
        try:
            png = vegalite_to_png(chart_spec, scale=2.0)
            s.shapes.add_picture(io.BytesIO(png), Inches(5.4), Inches(1.7),
                                 width=Inches(7.2))
        except Exception as e:
            self._para(self._box(s, 5.4, 3.0, 7.0, 1.0),
                       f"[chart unavailable: {e}]", 12, self.theme.ink_dim, first=True)

    def appendix_slide(self, items: list[dict]) -> None:
        """The audit trail: every number's SQL. The defensibility slide."""
        s = self._add()
        self._bg(s, self.theme.bg_dark)
        self._para(self._box(s, 0.8, 0.55, 11.7, 0.9), "Appendix — how these numbers were produced",
                   24, "FFFFFF", bold=True, font=self.theme.title_font, first=True)
        tf = self._box(s, 0.8, 1.7, 11.7, 5.2)
        first = True
        for it in items[:6]:
            self._para(tf, it.get("title", "query"), 14, self.theme.accent, bold=True,
                       font=self.theme.body_font, first=first)
            first = False
            sql = (it.get("sql", "") or "").strip()
            self._para(tf, sql or "(no SQL)", 11, "CADCDC", font="Courier New").space_after = Pt(10)

    # -- assemble ------------------------------------------------------------
    def build(self, deck: list[dict]) -> bytes:
        for slide in deck:
            t = slide.get("type")
            if t == "title":
                self.title_slide(slide.get("title", "Report"), slide.get("subtitle", ""))
            elif t == "summary":
                self.summary_slide(slide.get("title", "Summary"), slide.get("bullets", []))
            elif t == "chart":
                self.chart_slide(slide.get("title", ""), slide.get("takeaway", ""),
                                 slide.get("chart", {}))
            elif t == "appendix":
                self.appendix_slide(slide.get("items", []))
        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()


def build_deck(deck: list[dict], template_path: str | None = None) -> bytes:
    return DeckBuilder(template_path=template_path).build(deck)
